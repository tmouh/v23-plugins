"""
build-deck.py — CC-native NEW-BUILD generator for v23:om-builder.

Generates a Vanadium-grade .pptx directly from a pre-written content spec, by
cloning the house template (assets/v23-template.pptx — inherits the real master,
theme, fonts, color scheme, and the navy footer bar / page numbers) and rendering
the core DATA-DRIVEN layouts deterministically. This is the new-build counterpart
to apply-revision-edits.py: CC owns the file; CIP is reserved for visual polish
and bespoke-visual slides (full-bleed hero, annotated aerial, renderings, image
placement) the JS API can't do well.

COVERAGE (the deterministic core — ~80% of an OM):
  cover, kpi_strip, narrative, table, chart, two_column, section_divider.
NOT COVERED (route to CIP / manual): full-bleed photo heroes, annotated aerials,
maps, renderings, image placement, stacking plans, scatter/quadrant exhibits.
The generator marks any slide whose "layout" it doesn't implement as a labeled
placeholder slide so the deck is complete and CIP knows exactly what to finish.

House style + voice are NOT re-derived here — they live in design-system.md /
layout-system.md (visual) and voice-model.md + anti-ai-ruleset.md (writing voice).
CC must have already written the content in-voice before calling this; the
generator only places it.

─────────────────────────────────────────────────────────────────────────────────
CONTENT SPEC — CANONICAL JSON SCHEMA (Phase 3 Blueprint → Phase 4 Build contract)
─────────────────────────────────────────────────────────────────────────────────

The JSON spec is the interface between Phase 3 (Blueprint) and Phase 4 (Build).
Every Phase 3 blueprint must emit a spec that conforms to this schema exactly.
Phase 4 (this script) renders it without re-interpreting content decisions.

TOP-LEVEL KEYS:
  "template"   (string, optional) — absolute or relative path to v23-template.pptx.
               Defaults to ../assets/v23-template.pptx relative to this script file.
               Specs may omit this key; the default is always the skill's asset.
  "target"     (string, required) — output .pptx path.
  "deal_name"  (string, optional) — deal identifier used in notes/attribution.
  "slides"     (array, required) — ordered list of slide-spec objects.

SLIDE-SPEC OBJECT — keys by layout type:

  ALL layouts support:
    "layout"    (string, required) — one of: cover | kpi_strip | narrative |
                table | chart | two_column | section_divider. Any other value
                produces a labeled CIP-placeholder slide.
    "eyebrow"   (string, optional) — all-caps section/status label. On banded
                content slides it joins the band text as "EYEBROW  |  Title"
                (layout-system A-16 pattern); there is no standalone eyebrow
                shape on body slides (A-08). On covers it renders separately.
    "title"     (string, optional) — primary headline. Rendered as header-band
                text (Garamond 20pt bold white in the navy header band) on all
                content slides. Cover title uses a different position/size.
    "takeaway"  (string, optional) — italic italicized one-liner below the title.
                Shape name: "Takeaway" (semantic, per design-system override item 6).
    "source"    (string, optional) — citation line at slide bottom.
                Garamond 11pt italic (design-system §3 override item 8: confirmed
                Garamond 11pt from run-level XML inspection of seed deck 2026-06-11).
                Shape name: "Source" (semantic, per design-system override item 6).

  "cover":
    "title"     — deal name / address (Garamond 30pt bold, positioned below center)
    "eyebrow"   — sponsor | date line
    "subtitle"  — deal type / submarket context
    "tagline"   — equity raise / ask line
    NOTE: The cover hero background image (left-half full-bleed, x=0.0 y=0.0
    w=6.83 h=7.5) is a CIP/manual add. A note is written to slide notes.

  "kpi_strip":
    "kpis"  (array) — each item: {"value": "...", "label": "..."}
            Shape names: KPI_Val_N, KPI_Lbl_N (design-system override item 6).
            The KPI_Bg_N background rectangle is drawn as a filled rectangle.
            Up to 5 KPIs per strip recommended (design-system §A-07 observed).

  "narrative":
    "blocks" (array) — each item: {"lead": "...", "body": "..."}
             "lead" is rendered bold navy (the label before the em-dash).
             "body" supports **bold** inline markers: text wrapped in ** is
             rendered as a bold run; surrounding text is normal weight.
             Shape names: Block_1, Block_2, ... (semantic).

  "table":
    "table" — {"header": [...], "rows": [[...], ...]}
    Header row uses Garamond 11pt bold white on navy fill.
    Body rows use Garamond 10pt.

  "chart":
    "chart" — {"type": "line"|"column"|"bar"|"pie",
                "categories": [...],
                "series": [["series_name", [val, ...]], ...]}

  "two_column":
    "left"  — {"head": "...", "body": "..."}
    "right" — {"head": "...", "body": "..."}
    "body" supports **bold** inline markers (same as narrative).

  "section_divider":
    "title" — section name (ALL-CAPS, Garamond 32pt bold navy)

─────────────────────────────────────────────────────────────────────────────────
FULL EXAMPLE SPEC (copy-paste starting point for a new build):
─────────────────────────────────────────────────────────────────────────────────
{
  "template": "../assets/v23-template.pptx",
  "target": "/tmp/example-deck.pptx",
  "deal_name": "NPV Florida IOS Strategy",
  "slides": [
    {"layout":"cover",
     "eyebrow":"NORTH PARK VENTURES | MAY 2026",
     "title":"Industrial Outdoor Storage Strategy",
     "subtitle":"Programmatic Equity Partnership Opportunity — Florida",
     "tagline":"$50M Equity Raise | Sponsored by North Park Ventures"},
    {"layout":"kpi_strip","eyebrow":"EXECUTIVE SUMMARY",
     "title":"Transaction Overview",
     "takeaway":"Programmatic IOS acquisition strategy targeting central Florida.",
     "kpis":[{"value":"$50M","label":"LP EQUITY ASK"},
             {"value":"27%","label":"TARGET GROSS IRR"},
             {"value":"2.1x","label":"EQUITY MULTIPLE"},
             {"value":"6","label":"ASSETS IN PIPELINE"},
             {"value":"18 MO","label":"DEPLOYMENT HORIZON"}],
     "source":"Source: Newmark, 2025."},
    {"layout":"narrative","eyebrow":"INVESTMENT THESIS",
     "title":"Why Industrial Outdoor Storage, Why Now",
     "takeaway":"IOS represents the highest-conviction value-add play in Florida industrial.",
     "blocks":[
       {"lead":"THE OPPORTUNITY",
        "body":"Florida IOS vacancy has compressed to **2.1%** — a 15-year low — while new supply is constrained by **zoning friction** across all major MSAs."},
       {"lead":"THE STRATEGY",
        "body":"Acquire **6–8 stabilized IOS sites** across the I-4 corridor at **$55–75/SF**, execute light-touch improvement, and refinance within **18 months** of close."}
     ],
     "source":"Source: CoStar Industrial, Q1 2026."},
    {"layout":"table","eyebrow":"PIPELINE","title":"Live Deployable Deals",
     "takeaway":"Six assets identified; three under LOI as of June 2026.",
     "table":{"header":["Asset","Market","Size (SF)","Ask","Status"],
              "rows":[["Skyway Circle","Melbourne, FL","42,000","$2.9M","Under LOI"],
                      ["Gatlin Yard","Orlando, FL","58,000","$4.1M","Under LOI"],
                      ["Lakeland IOS","Lakeland, FL","35,000","$2.3M","Tracking"]]},
     "source":"Source: North Park Ventures pipeline, June 2026."},
    {"layout":"chart","eyebrow":"MARKET","title":"Florida IOS Vacancy Rate 2019–2026",
     "takeaway":"Vacancy has compressed 640 bps over six years with no supply relief.",
     "chart":{"type":"line","categories":["2019","2020","2021","2022","2023","2024","2025","2026"],
              "series":[["Vacancy %",[8.7,7.9,6.2,4.8,3.5,2.9,2.3,2.1]]]},
     "source":"Source: CoStar Industrial, Q1 2026."},
    {"layout":"two_column","eyebrow":"SPONSOR","title":"North Park Ventures Track Record",
     "takeaway":"18 acquisitions, zero losses, $380M in total capitalization since 2019.",
     "left":{"head":"PLATFORM OVERVIEW",
             "body":"North Park Ventures is a **Florida-focused** industrial operator with **18 closed acquisitions** across the I-4 and I-75 corridors since 2019."},
     "right":{"head":"SELECTED REALIZATIONS",
              "body":"Skyway Circle (Dec 2024): **27% gross IRR**, **2.1x EM** on a 14-month hold. Gatlin Phase I (Mar 2024): **31% gross IRR**, **1.9x EM** on a 12-month hold."},
     "source":"Source: North Park Ventures, as of June 2026."},
    {"layout":"section_divider","title":"SPONSOR TRACK RECORD"}
  ]
}

─────────────────────────────────────────────────────────────────────────────────
XML-CLONE ESCAPE HATCH (SANCTIONED FALLBACK — DO NOT REMOVE)
─────────────────────────────────────────────────────────────────────────────────

For arrangements python-pptx cannot express (complex grouped shapes, annotated
aerials, stacking plans, certain photo-grid constructions, the A-04 2×2 block grid
with named Block_* shapes, the A-16 bottom KPI strip at exact positions), the
production-proven method is:

  1. UNPACK:  unzip <source_template.pptx> -d _build/unpacked
  2. CLONE:   copy ppt/slides/slideN.xml  → ppt/slides/slideM.xml
              copy ppt/slides/_rels/slideN.xml.rels → ppt/slides/_rels/slideM.xml.rels
  3. FIX rIds: all rId values in the cloned .xml.rels must be unique (rename r1→rM1 etc.)
              and the cloned slide's rId references must match.
  4. REGISTER: add the new slide to ppt/presentation.xml sldIdLst (new unique id attr)
              and add a corresponding entry in ppt/_rels/presentation.xml.rels.
  5. REPACK:  cd _build/unpacked && zip -r ../../output.pptx . -x "*.DS_Store"

This procedure was used successfully in:
  - apply-revision-edits.py (all B-chain revision cycles)
  - D-NPV endgame build (June 2026) — the A-04 exec summary 2×2 grid and
    A-16 realized case study KPI strip were cloned from the seed deck, not
    generated by python-pptx.

When to use the clone path instead of this script:
  - Any slide type in CIP_LAYOUTS (full_bleed, aerial, map, rendering, etc.)
  - A-04 / A-16 / A-15 layouts with named positioned shapes that must match
    design-system geometry exactly (python-pptx cannot reproduce named grouped rects)
  - Any slide where the visual result, not just the text content, must be
    pixel-faithful to the template

Never build a large new python-pptx subsystem to replace the clone path.
The clone path is correct, minimal, and proven. Extend it, don't replace it.
"""

import argparse
import json
import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn

# ── House style constants ──────────────────────────────────────────────────────
# Values verified from design-system.md and layout-system.md observed geometry.
FONT = "Garamond"
NAVY = "1F3A5F"       # D-01 / A-08: header band fill, primary navy
NAVY_DARK = "0F2540"
MIDBLUE = "5B7FA8"
PANEL = "E8EEF6"
BODY = "485269"
MUTED = "9AA3B2"
WHITE = "FFFFFF"      # D-04: header text color (scheme:bg1 / #FFFFFF)

# ── Canvas (design-system D-12: 13.33" × 7.50" widescreen, all V23 equity OMs)
CANVAS_W = 13.33
CANVAS_H = 7.50

# ── Header-band shell (layout-system A-08: exact observed geometry) ───────────
# Rectangle: x=0.0 y=0.29 w=13.33 h=0.50 fill #1F3A5F
BAND_X, BAND_Y, BAND_W, BAND_H = 0.0, 0.29, 13.33, 0.50
# Header text TextBox: x=0.41 y=0.33 w=12.52 h=0.44, Garamond 20pt bold white
BAND_TXT_X, BAND_TXT_Y, BAND_TXT_W, BAND_TXT_H = 0.41, 0.33, 12.52, 0.44
BAND_FONT_SIZE = 20  # design-system D-04: Garamond 20pt bold white

# ── Content area geometry (relative to band shell) ───────────────────────────
# These positions are used for text-box placement BELOW the header band.
MARGIN_L = 0.83
USABLE_W = 11.67          # 0.83 .. 12.50

# layout-system A-08 content area begins y≈1.03–1.06. There is no standalone
# eyebrow shape on banded content slides — the eyebrow joins the band text
# (layout-system A-16: "REALIZED CASE STUDY | 3161 SKYWAY CIRCLE, ...").
TAKEAWAY_Y, TAKEAWAY_H = 1.33, 0.32
BODY_TOP = 1.94
SOURCE_Y, SOURCE_H = 6.97, 0.17

CHART_TYPES = {
    "line": XL_CHART_TYPE.LINE,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "pie": XL_CHART_TYPE.PIE,
}

# ── Inline-bold marker pattern ─────────────────────────────────────────────────
_BOLD_RE = re.compile(r"\*\*(.+?)\*\*", re.DOTALL)


def _hex(c):
    return RGBColor.from_string(c)


def _box(slide, x, y, w, h, name=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    if name:
        tb.name = name
    return tb


def _set(tb, text, size, color=BODY, bold=False, italic=False, caps=False,
         align=PP_ALIGN.LEFT, font=FONT):
    """Set text on a textbox with uniform formatting across all lines."""
    tf = tb.text_frame
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line.upper() if caps else line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = _hex(color)
    return tb


def _set_inline_bold(tb, text, size, color=BODY, base_bold=False,
                     italic=False, align=PP_ALIGN.LEFT, font=FONT):
    """Set text supporting **bold** inline markers.

    Splits the text into segments on ** markers, emitting alternating
    normal/bold runs within a single paragraph. Newlines produce new
    paragraphs (each inheriting the same base formatting).
    """
    tf = tb.text_frame
    lines = text.split("\n")
    for li, line in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = align
        # Split on **marker** boundaries
        segments = _BOLD_RE.split(line)
        # _BOLD_RE.split yields: [before, bold_content, after, bold_content2, ...]
        # Odd-indexed items (1, 3, …) are the captured bold groups.
        for si, seg in enumerate(segments):
            if not seg:
                continue
            r = p.add_run()
            r.text = seg
            f = r.font
            f.name = font
            f.size = Pt(size)
            f.italic = italic
            f.color.rgb = _hex(color)
            # Odd indices are captured bold groups
            f.bold = True if (si % 2 == 1) else base_bold
    return tb


def _draw_header_band(slide):
    """Draw the navy header band on a content slide.

    layout-system A-08 exact geometry: Rectangle x=0.0 y=0.29 w=13.33 h=0.50
    fill #1F3A5F. This is the permanent template fixture on ALL content slides
    in the Garamond system. We draw it here to ensure programmatically-generated
    slides carry it even when the template's master doesn't inject it into each
    slide XML directly.
    """
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE (value 1)
        Inches(BAND_X), Inches(BAND_Y), Inches(BAND_W), Inches(BAND_H)
    )
    shape.name = "Section_Header"  # design-system override item 6: semantic name
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(NAVY)  # design-system D-01: fill #1F3A5F
    shape.line.fill.background()            # no border


def _header(slide, s):
    """Title-in-band + takeaway — the standard top-of-slide block.

    The title is rendered in the header band position per design-system D-04:
    TextBox x=0.41 y=0.33 w=12.52 h=0.44, Garamond 20pt bold white.
    No standalone eyebrow shape exists on banded slides (layout-system A-08);
    when an eyebrow is given it joins the band text with a pipe, per the
    observed A-16 pattern ("REALIZED CASE STUDY | 3161 SKYWAY CIRCLE, ...").
    """
    _draw_header_band(slide)
    band_text = s.get("title", "")
    if s.get("eyebrow"):
        band_text = (s["eyebrow"].upper() + "  |  " + band_text) if band_text \
            else s["eyebrow"].upper()
    if band_text:
        # design-system D-04: Garamond 20pt bold white at band text position
        _set(_box(slide, BAND_TXT_X, BAND_TXT_Y, BAND_TXT_W, BAND_TXT_H, "Title"),
             band_text, BAND_FONT_SIZE, WHITE, bold=True)
    if s.get("takeaway"):
        # design-system override item 6: semantic name "Takeaway"
        _set(_box(slide, MARGIN_L, TAKEAWAY_Y, USABLE_W, TAKEAWAY_H, "Takeaway"),
             s["takeaway"], 13, BODY, italic=True)


def _source(slide, s):
    """Source line at slide bottom.

    design-system §3 override item 8 (resolved 2026-06-11): Garamond 11pt.
    Run-level XML inspection of seed deck confirmed slide 15 = Garamond 11pt,
    slide 28 = Garamond 12pt. New builds use 11pt as the canonical value.
    Shape name: "Source" (design-system override item 6).
    """
    if s.get("source"):
        _set(_box(slide, MARGIN_L, SOURCE_Y, USABLE_W, SOURCE_H, "Source"),
             s["source"], 11, MUTED, italic=True)  # 11pt per design-system §3 override item 8


# ── Layout renderers ────────────────────────────────────────────────────────
def render_cover(slide, s):
    """Cover slide — text-only positions. Hero image is CIP/manual.

    Note: The left-half full-bleed photo (x=0.0 y=0.0 w=6.83 h=7.5) is placed
    by CIP or manually; see XML-clone escape hatch in module docstring.
    """
    _set(_box(slide, MARGIN_L, 4.69, 7.5, 0.28, "CoverEyebrow"),
         s.get("eyebrow", ""), 11, NAVY, bold=True, caps=True)
    _set(_box(slide, MARGIN_L, 4.95, 7.5, 1.38, "CoverTitle"),
         s.get("title", ""), 30, NAVY, bold=True)
    _set(_box(slide, MARGIN_L, 6.35, 7.5, 0.35, "CoverSubtitle"),
         s.get("subtitle", ""), 14, BODY)
    if s.get("tagline"):
        _set(_box(slide, MARGIN_L, 6.69, 7.5, 0.28, "CoverTagline"),
             s["tagline"], 11, MIDBLUE, bold=True)
    slide.notes_slide.notes_text_frame.text = (
        "[CIP/manual: place full-bleed hero image behind the cover title block. "
        "A-01 geometry: left-half photo x=0.0 y=0.0 w=6.83 h=7.5 (design-system D-06 / layout-system A-01).]"
    )


def render_kpi_strip(slide, s):
    """KPI strip with named background rects.

    Shape naming per design-system override item 6:
      KPI_Bg_N  — filled rectangle background
      KPI_Val_N — large value text
      KPI_Lbl_N — ALL-CAPS label text
    """
    _header(slide, s)
    kpis = s.get("kpis", [])
    n = max(len(kpis), 1)
    gap = 0.13
    tile_w = (USABLE_W - (n - 1) * gap) / n
    for i, k in enumerate(kpis):
        x = MARGIN_L + i * (tile_w + gap)
        # Background rect — design-system override item 6: KPI_Bg_N
        bg = slide.shapes.add_shape(
            1,  # rectangle
            Inches(x), Inches(BODY_TOP), Inches(tile_w), Inches(1.1)
        )
        bg.name = f"KPI_Bg_{i+1}"
        bg.fill.solid()
        bg.fill.fore_color.rgb = _hex(PANEL)
        bg.line.fill.background()
        # Value — design-system override item 6: KPI_Val_N
        _set(_box(slide, x, BODY_TOP + 0.1, tile_w, 0.5, f"KPI_Val_{i+1}"),
             k.get("value", ""), 20, NAVY, bold=True, align=PP_ALIGN.CENTER)
        # Label — design-system override item 6: KPI_Lbl_N
        _set(_box(slide, x, BODY_TOP + 0.62, tile_w, 0.4, f"KPI_Lbl_{i+1}"),
             k.get("label", ""), 9, MUTED, caps=True, align=PP_ALIGN.CENTER)
    _source(slide, s)


def render_narrative(slide, s):
    """Narrative prose slide.

    Block body text supports **bold** inline markers — split into bold/normal
    runs within the same paragraph. Shape names: Block_1, Block_2, ...
    (design-system override item 6: semantic naming).
    """
    _header(slide, s)
    blocks = s.get("blocks", [])
    y = BODY_TOP
    block_h = min(1.5, (6.6 - BODY_TOP) / max(len(blocks), 1))
    for i, b in enumerate(blocks):
        tb = _box(slide, MARGIN_L, y, USABLE_W, block_h, f"Block_{i+1}")
        tf = tb.text_frame
        p = tf.paragraphs[0]
        # Lead (bold navy label before the em-dash)
        if b.get("lead"):
            r = p.add_run()
            r.text = b["lead"] + " — "
            r.font.name = FONT
            r.font.size = Pt(13)
            r.font.bold = True
            r.font.color.rgb = _hex(NAVY)
        # Body with **bold** inline marker support
        body_text = b.get("body", "")
        if body_text:
            _append_inline_bold_runs(p, body_text, size_pt=13, color=BODY,
                                     base_bold=False, font=FONT)
        y += block_h
    _source(slide, s)


def _append_inline_bold_runs(paragraph, text, size_pt, color, base_bold=False,
                              font=FONT):
    """Append runs with **bold** inline marker support to an existing paragraph.

    Segments on ** boundaries: odd-indexed captures are bold, even-indexed are
    normal weight. This is the canonical inline-bold implementation — used by
    render_narrative and render_two_column.
    """
    segments = _BOLD_RE.split(text)
    for si, seg in enumerate(segments):
        if not seg:
            continue
        r = paragraph.add_run()
        r.text = seg
        r.font.name = font
        r.font.size = Pt(size_pt)
        r.font.color.rgb = _hex(color)
        r.font.bold = True if (si % 2 == 1) else base_bold


def render_table(slide, s):
    _header(slide, s)
    t = s.get("table", {})
    header = t.get("header", [])
    rows = t.get("rows", [])
    n_rows = len(rows) + (1 if header else 0)
    n_cols = max(len(header), max((len(r) for r in rows), default=0))
    if n_rows == 0 or n_cols == 0:
        return
    gframe = slide.shapes.add_table(n_rows, n_cols, Inches(MARGIN_L), Inches(BODY_TOP),
                                    Inches(USABLE_W), Inches(0.35 * n_rows))
    tbl = gframe.table
    r0 = 0
    if header:
        for c, htext in enumerate(header):
            cell = tbl.cell(0, c)
            cell.text = str(htext)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = FONT
                    run.font.size = Pt(11)  # design-system D-04: Garamond header weight
                    run.font.bold = True
                    run.font.color.rgb = _hex(WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _hex(NAVY)
        r0 = 1
    for ri, row in enumerate(rows):
        for c in range(n_cols):
            cell = tbl.cell(ri + r0, c)
            cell.text = str(row[c]) if c < len(row) else ""
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = FONT
                    run.font.size = Pt(10)
                    run.font.color.rgb = _hex(BODY)
    _source(slide, s)


def render_chart(slide, s):
    _header(slide, s)
    c = s.get("chart", {})
    data = CategoryChartData()
    data.categories = c.get("categories", [])
    for name, vals in c.get("series", []):
        data.add_series(name, tuple(vals))
    ctype = CHART_TYPES.get(c.get("type", "line"), XL_CHART_TYPE.LINE)
    slide.shapes.add_chart(ctype, Inches(MARGIN_L), Inches(BODY_TOP),
                           Inches(USABLE_W), Inches(3.6), data)
    _source(slide, s)


def render_two_column(slide, s):
    """Two-column layout. Body text in each column supports **bold** inline markers."""
    _header(slide, s)
    half = (USABLE_W - 0.33) / 2
    for side, x in (("left", MARGIN_L), ("right", MARGIN_L + half + 0.33)):
        col = s.get(side, {})
        if col.get("head"):
            _set(_box(slide, x, BODY_TOP, half, 0.3, f"{side}_head"),
                 col["head"], 13, NAVY, bold=True)
        if col.get("body"):
            tb = _box(slide, x, BODY_TOP + 0.35, half, 4.0, f"{side}_body")
            p = tb.text_frame.paragraphs[0]
            _append_inline_bold_runs(p, col["body"], size_pt=13, color=BODY)
    _source(slide, s)


def render_section_divider(slide, s):
    _set(_box(slide, MARGIN_L, 3.0, USABLE_W, 1.0, "DividerTitle"),
         s.get("title", ""), 32, NAVY, bold=True, caps=True)


def render_placeholder(slide, s):
    """A slide whose layout the generator doesn't implement — leave a clear,
    labeled placeholder so the deck is complete and CIP knows what to finish."""
    _header(slide, s)
    label = f"[ CIP / MANUAL: {s.get('layout','?')} ]  {s.get('note','')}"
    _set(_box(slide, MARGIN_L, 3.0, USABLE_W, 1.0, "Placeholder"),
         label, 14, MIDBLUE, bold=True, align=PP_ALIGN.CENTER)
    _source(slide, s)


RENDERERS = {
    "cover": render_cover,
    "kpi_strip": render_kpi_strip,
    "narrative": render_narrative,
    "table": render_table,
    "chart": render_chart,
    "two_column": render_two_column,
    "section_divider": render_section_divider,
}

# Layouts the generator does NOT implement → placeholder + flagged for CIP.
# These must be handled via the XML-clone escape hatch (see module docstring).
CIP_LAYOUTS = {"full_bleed", "aerial", "map", "rendering", "photo_grid",
               "stacking_plan", "scatter", "quadrant", "annotated_aerial"}


def build(spec):
    here = os.path.dirname(os.path.abspath(__file__))
    template = spec.get("template") or os.path.join(here, "..", "assets", "v23-template.pptx")
    template = os.path.abspath(template)
    if not os.path.isfile(template):
        raise SystemExit(f"ERROR: template not found: {template} (run make-template.py first)")
    prs = Presentation(template)
    # Find the Blank layout (house pattern: manual shapes on Blank)
    blank = None
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == "Blank":
                blank = layout
                break
    if blank is None:
        blank = prs.slide_masters[0].slide_layouts[-1]

    report = []
    for i, s in enumerate(spec["slides"], 1):
        slide = prs.slides.add_slide(blank)
        lay = s.get("layout", "narrative")
        if lay in RENDERERS:
            RENDERERS[lay](slide, s)
            report.append(f"slide {i}: {lay} (rendered)")
        elif lay in CIP_LAYOUTS:
            render_placeholder(slide, s)
            report.append(f"slide {i}: {lay} -> PLACEHOLDER (route to CIP)")
        else:
            render_placeholder(slide, s)
            report.append(f"slide {i}: UNKNOWN layout '{lay}' -> PLACEHOLDER")
    prs.save(spec["target"])
    return report


def main(argv=None):
    p = argparse.ArgumentParser(description="Generate a Vanadium-grade .pptx from a content spec.")
    p.add_argument("spec", help="Path to JSON content spec.")
    args = p.parse_args(argv)
    with open(args.spec, encoding="utf-8") as f:
        spec = json.load(f)
    report = build(spec)
    print("\n".join(report))
    print(f"\nGenerated {len(report)} slides -> {spec['target']}")
    print("NOTE: structure/text/geometry are deterministic; open the file to judge VISUAL quality.")
    print("Placeholder slides are flagged above — finish those in PowerPoint/CIP.")


if __name__ == "__main__":
    main()
