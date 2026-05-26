"""
apply-revision-edits.py — CC-native revision engine for v23:om-builder.

Applies a deterministic edit spec directly to a COPY of a .pptx via python-pptx,
preserving run-level formatting, verifying preconditions before each edit, and
writing a finished revised file plus a change report. This is the fast,
quota-free, lock-safe alternative to hand-executing a CIP edit script for the
deterministic majority of revision work (text, numbers, table cells, chart data).

Design contract (mirrors the CIP edit-script discipline):
  - Never touches the input file. Always copies to the output path first.
  - Targets every shape by numeric shape_id (the same IDs the inventory emits).
  - Verifies a precondition (expected current text substring) before editing.
    On mismatch it records a FAILURE and SKIPS that edit — it does not guess.
  - Preserves the shape's existing run formatting (font, size, color, bold,
    italic) when replacing text.
  - Prints a per-edit report (PASS / SKIP-precondition / ERROR) and exits
    non-zero if any edit failed, so CC can halt and surface it.

Edit spec (JSON): {"target": "<output.pptx>", "source": "<input.pptx>",
                   "edits": [ {op...}, ... ]}
Ops:
  {"op":"set_text","slide":2,"shape_id":15,"expect":"27.0%","text":"20%+"}
  {"op":"set_rich_text","slide":2,"shape_id":7,"expect":"The Opportunity",
      "segments":[["THE OPPORTUNITY ",true],["Industrial Outdoor Storage...",false]]}
  {"op":"set_cell","slide":8,"shape_id":6,"row":0,"col":0,"expect":"Submarket","text":"Submarket (FL)"}
  {"op":"set_chart_series","slide":5,"shape_id":201,
      "categories":["2020",...],"series":[["IOS",[100,...]],["Bulk Warehouse",[...]]]}
  {"op":"delete_slide","slide":20,"expect_title":"SPONSOR"}

`slide` is 1-based (matches how humans/inventory refer to slides).

Usage:
  python apply-revision-edits.py edits.json
  python apply-revision-edits.py --source in.pptx --target out.pptx --inline '<json>'
"""

import argparse
import copy
import json
import os
import shutil
import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def _capture_font(run):
    """Snapshot a run's formatting so it can be re-applied after a text reset."""
    f = run.font
    snap = {"name": f.name, "size": f.size, "bold": f.bold, "italic": f.italic, "rgb": None}
    try:
        if f.color is not None and f.color.type is not None:
            snap["rgb"] = f.color.rgb
    except Exception:
        pass
    return snap


def _apply_font(run, snap):
    f = run.font
    if snap.get("name") is not None:
        f.name = snap["name"]
    if snap.get("size") is not None:
        f.size = snap["size"]
    if snap.get("bold") is not None:
        f.bold = snap["bold"]
    if snap.get("italic") is not None:
        f.italic = snap["italic"]
    if snap.get("rgb") is not None:
        f.color.rgb = snap["rgb"]


def _first_run(text_frame):
    for p in text_frame.paragraphs:
        if p.runs:
            return p.runs[0]
    return None


def find_shape(slide, shape_id):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    return None


def set_shape_text(shape, new_text):
    """Replace all text in a shape with new_text (may contain \\n for multiple
    paragraphs), inheriting the shape's dominant run formatting. Flattens inline
    formatting variation — use set_rich_text for shapes with inline bold."""
    tf = shape.text_frame
    template = _first_run(tf)
    snap = _capture_font(template) if template is not None else None
    lines = new_text.split("\n")
    tf.clear()  # leaves one empty paragraph
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        if snap:
            _apply_font(run, snap)


def set_rich_text(shape, segments):
    """Replace text with a list of [text, bold] segments in one paragraph,
    inheriting font name/size/color from the shape's first run, overriding bold
    per segment. For multi-paragraph rich text, pass segments with a literal
    '\\n' text segment to break lines."""
    tf = shape.text_frame
    template = _first_run(tf)
    snap = _capture_font(template) if template is not None else None
    tf.clear()
    para = tf.paragraphs[0]
    for seg in segments:
        text, bold = seg[0], (seg[1] if len(seg) > 1 else False)
        for j, line in enumerate(text.split("\n")):
            if j > 0:
                para = tf.add_paragraph()
            run = para.add_run()
            run.text = line
            if snap:
                _apply_font(run, snap)
            run.font.bold = bool(bold)


def set_cell_text(table, row, col, new_text):
    cell = table.cell(row, col)
    tf = cell.text_frame
    template = _first_run(tf)
    snap = _capture_font(template) if template is not None else None
    tf.clear()
    run = tf.paragraphs[0].add_run()
    run.text = new_text
    if snap:
        _apply_font(run, snap)


def set_chart_series(chart, categories, series):
    data = CategoryChartData()
    data.categories = categories
    for name, vals in series:
        data.add_series(name, tuple(vals))
    chart.replace_data(data)


_RID_ATTR = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


def delete_slide(prs, slide):
    """Remove a slide FULLY: drop its relationship (so the slide part is not
    serialized) AND remove the sldId entry. Removing only the sldId leaves an
    orphaned slideN.xml part that collides when new slides are later added.
    Slide order in _sldIdLst matches prs.slides, so index alignment is safe."""
    target_idx = list(prs.slides).index(slide)
    xml_slides = prs.slides._sldIdLst
    sld = list(xml_slides)[target_idx]
    rId = sld.get(_RID_ATTR)
    if rId:
        prs.part.drop_rel(rId)
    xml_slides.remove(sld)


def _find_layout(prs, layout_name):
    """Find a slide layout by its name across all slide masters."""
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == layout_name:
                return layout
    # fallback: first master's first layout matching loosely
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout_name.lower() in (layout.name or "").lower():
                return layout
    return None


def _add_textbox(slide, spec):
    """Add a text box from a spec dict: {text, x_in, y_in, w_in, h_in,
    font, size_pt, color (hex no #), bold, italic, align}."""
    tb = slide.shapes.add_textbox(
        Inches(spec["x_in"]), Inches(spec["y_in"]),
        Inches(spec["w_in"]), Inches(spec["h_in"]),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    lines = spec["text"].split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        f = run.font
        if spec.get("font"):
            f.name = spec["font"]
        if spec.get("size_pt"):
            f.size = Pt(spec["size_pt"])
        if spec.get("bold") is not None:
            f.bold = bool(spec["bold"])
        if spec.get("italic") is not None:
            f.italic = bool(spec["italic"])
        if spec.get("color"):
            f.color.rgb = RGBColor.from_string(spec["color"].lstrip("#"))
    if spec.get("name"):
        tb.name = spec["name"]
    return tb


def insert_slide(prs, at_index, layout_name, textboxes):
    """Add a slide on the named layout, move it to at_index (1-based), and
    populate it with the supplied text boxes. Returns the new slide.
    NOTE: only supports text boxes — images/charts/tables on a NEW slide are
    out of scope for the engine (route those to CIP/manual)."""
    layout = _find_layout(prs, layout_name)
    if layout is None:
        raise ValueError(f"layout '{layout_name}' not found")
    slide = prs.slides.add_slide(layout)  # appended at end
    for spec in (textboxes or []):
        _add_textbox(slide, spec)
    # move from end to the requested 0-based position
    target_idx = at_index - 1
    xml_slides = prs.slides._sldIdLst
    sld_ids = list(xml_slides)
    new_sld = sld_ids[-1]
    xml_slides.remove(new_sld)
    xml_slides.insert(target_idx, new_sld)
    return slide


def set_geometry(shape, spec):
    """Resize/move a shape. Any of x_in/y_in/w_in/h_in may be supplied."""
    if spec.get("x_in") is not None:
        shape.left = Inches(spec["x_in"])
    if spec.get("y_in") is not None:
        shape.top = Inches(spec["y_in"])
    if spec.get("w_in") is not None:
        shape.width = Inches(spec["w_in"])
    if spec.get("h_in") is not None:
        shape.height = Inches(spec["h_in"])


def shape_text(shape):
    if shape.has_text_frame:
        return shape.text_frame.text
    return ""


def run(spec, report):
    source = spec["source"]
    target = spec["target"]
    if os.path.abspath(source) == os.path.abspath(target):
        raise SystemExit("ERROR: source and target must differ (never edit the original in place).")
    shutil.copy2(source, target)
    prs = Presentation(target)
    n_slides = len(prs.slides)
    failures = 0

    # Order: content edits first (stable indices), then deletes (descending
    # slide order), then inserts (ascending at_index, interpreted against the
    # post-delete deck). One structural change per run is the simplest safe case.
    structural = ("delete_slide", "insert_slide")
    edits = [e for e in spec["edits"] if e["op"] not in structural]
    deletes = sorted([e for e in spec["edits"] if e["op"] == "delete_slide"],
                     key=lambda e: e["slide"], reverse=True)
    inserts = sorted([e for e in spec["edits"] if e["op"] == "insert_slide"],
                     key=lambda e: e["at_index"])

    for e in edits:
        op = e["op"]
        si = e["slide"] - 1
        if si < 0 or si >= n_slides:
            report.append(f"ERROR  {op} slide {e['slide']}: out of range (deck has {n_slides})")
            failures += 1
            continue
        slide = prs.slides[si]
        try:
            if op in ("set_text", "set_rich_text"):
                sh = find_shape(slide, e["shape_id"])
                if sh is None:
                    report.append(f"ERROR  {op} s{e['slide']} id{e['shape_id']}: shape not found")
                    failures += 1
                    continue
                cur = shape_text(sh)
                if e.get("expect") and e["expect"] not in cur:
                    report.append(f"SKIP   {op} s{e['slide']} id{e['shape_id']}: precondition '{e['expect']}' "
                                  f"not in current text '{cur[:60]}'")
                    failures += 1
                    continue
                if op == "set_text":
                    set_shape_text(sh, e["text"])
                    report.append(f"PASS   set_text s{e['slide']} id{e['shape_id']}: -> '{e['text'][:60]}'")
                else:
                    set_rich_text(sh, e["segments"])
                    report.append(f"PASS   set_rich_text s{e['slide']} id{e['shape_id']}")
            elif op == "set_cell":
                sh = find_shape(slide, e["shape_id"])
                if sh is None or not sh.has_table:
                    report.append(f"ERROR  set_cell s{e['slide']} id{e['shape_id']}: not a table")
                    failures += 1
                    continue
                cell = sh.table.cell(e["row"], e["col"])
                if e.get("expect") and e["expect"] not in cell.text:
                    report.append(f"SKIP   set_cell s{e['slide']} id{e['shape_id']} ({e['row']},{e['col']}): "
                                  f"precondition '{e['expect']}' not in '{cell.text[:40]}'")
                    failures += 1
                    continue
                set_cell_text(sh.table, e["row"], e["col"], e["text"])
                report.append(f"PASS   set_cell s{e['slide']} id{e['shape_id']} ({e['row']},{e['col']}): -> '{e['text']}'")
            elif op == "set_chart_series":
                sh = find_shape(slide, e["shape_id"])
                if sh is None or not sh.has_chart:
                    report.append(f"ERROR  set_chart_series s{e['slide']} id{e['shape_id']}: not a chart")
                    failures += 1
                    continue
                set_chart_series(sh.chart, e["categories"], e["series"])
                report.append(f"PASS   set_chart_series s{e['slide']} id{e['shape_id']}: {len(e['series'])} series")
            elif op == "set_geometry":
                sh = find_shape(slide, e["shape_id"])
                if sh is None:
                    report.append(f"ERROR  set_geometry s{e['slide']} id{e['shape_id']}: shape not found")
                    failures += 1
                    continue
                set_geometry(sh, e)
                report.append(f"PASS   set_geometry s{e['slide']} id{e['shape_id']}")
            else:
                report.append(f"ERROR  unknown op '{op}'")
                failures += 1
        except Exception as ex:
            report.append(f"ERROR  {op} s{e['slide']} id{e.get('shape_id')}: {type(ex).__name__}: {ex}")
            failures += 1

    for e in deletes:
        si = e["slide"] - 1
        slide = prs.slides[si]
        title = ""
        if slide.shapes.title is not None:
            title = slide.shapes.title.text
        if e.get("expect_title") and e["expect_title"] not in title:
            # also scan all shapes for the expected title text
            found = any(e["expect_title"] in shape_text(sh) for sh in slide.shapes)
            if not found:
                report.append(f"SKIP   delete_slide {e['slide']}: title precondition '{e['expect_title']}' not found")
                failures += 1
                continue
        delete_slide(prs, slide)
        report.append(f"PASS   delete_slide {e['slide']} (was: '{e.get('expect_title','')}')")

    for e in inserts:
        try:
            insert_slide(prs, e["at_index"], e["layout"], e.get("textboxes", []))
            report.append(f"PASS   insert_slide at {e['at_index']} (layout '{e['layout']}', "
                          f"{len(e.get('textboxes', []))} text boxes)")
        except Exception as ex:
            report.append(f"ERROR  insert_slide at {e.get('at_index')}: {type(ex).__name__}: {ex}")
            failures += 1

    prs.save(target)
    return failures


def main(argv=None):
    p = argparse.ArgumentParser(description="Apply a deterministic edit spec to a COPY of a .pptx (CC-native revision engine).")
    p.add_argument("spec", nargs="?", help="Path to JSON edit spec.")
    p.add_argument("--source", help="Source .pptx (overrides spec).")
    p.add_argument("--target", help="Output .pptx (overrides spec).")
    p.add_argument("--inline", help="Inline JSON spec string (instead of a file).")
    args = p.parse_args(argv)

    if args.inline:
        spec = json.loads(args.inline)
    elif args.spec:
        with open(args.spec, encoding="utf-8") as f:
            spec = json.load(f)
    else:
        raise SystemExit("Provide a spec file or --inline JSON.")
    if args.source:
        spec["source"] = args.source
    if args.target:
        spec["target"] = args.target

    report = []
    failures = run(spec, report)
    print("\n".join(report))
    print(f"\n{'='*60}")
    print(f"Applied {len(spec['edits'])} edits | {failures} failed/skipped")
    print(f"Output: {spec['target']}")
    if failures:
        print("FAILURES PRESENT — review SKIP/ERROR lines above. Output written but incomplete.")
        sys.exit(1)
    print("ALL EDITS APPLIED CLEANLY.")


if __name__ == "__main__":
    main()
