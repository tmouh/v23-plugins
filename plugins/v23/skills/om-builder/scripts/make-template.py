"""
make-template.py — derive a blank Vanadium house-style template from a reference deck.

The new-build generator (build-deck.py) clones a template .pptx so it inherits the
real Vanadium slide master, theme (fonts + color scheme), and slide layouts — the
single biggest fidelity lever, and free. This script produces that template by
taking a finished production deck and stripping its content slides, leaving only
the master / layouts / theme.

Usage:
  python make-template.py <reference.pptx> <out-template.pptx>

Run once from any house-style production deck (e.g. a copy of the NPV v2 deck or
105 N 13th). Re-run if the house master/theme changes. Work on a COPY of the
reference if the original may be open in PowerPoint (file lock).
"""

import sys
from pptx import Presentation


RID_ATTR = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


def make_template(reference_path, out_path):
    prs = Presentation(reference_path)
    n_before = len(prs.slides)
    # Remove every content slide FULLY: drop the relationship (so the slide part
    # is unreachable and not serialized) AND remove the sldId entry. Removing only
    # the sldId leaves orphaned slideN.xml parts that collide when new slides are
    # later added (duplicate-part corruption). The master/layouts/theme remain.
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        rId = sld.get(RID_ATTR)
        if rId:
            prs.part.drop_rel(rId)
        xml_slides.remove(sld)
    prs.save(out_path)
    # Report the layouts the generator can target.
    chk = Presentation(out_path)
    layouts = []
    for master in chk.slide_masters:
        for layout in master.slide_layouts:
            layouts.append(layout.name)
    print(f"Stripped {n_before} content slides.")
    print(f"Template saved: {out_path}")
    print(f"Available layouts ({len(layouts)}): {', '.join(layouts)}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        raise SystemExit("usage: python make-template.py <reference.pptx> <out-template.pptx>")
    make_template(sys.argv[1], sys.argv[2])
