# -*- coding: utf-8 -*-
"""
verify-deck.py — Phase 5b verification harness for v23:om-builder.

Codifies the ledger re-check as a deterministic script: every critical figure
from 00-DATA-AUDIT-TRAIL.md must be PRESENT in the deck text, and every
blueprint §8 stop-list / never-reintroduce token (plus wrong-deck remnants from
the clone seed) must be ABSENT. First used on the Coastal GA build (2026-06-11):
70/70 present, 0 violations.

USAGE
    python verify-deck.py <deck.pptx> <checks.py>

The checks file is a deal-local python file (lives in x V23\\_build\\) defining:

    MUST = [
        # (label, regex, audit_trail_id) — figure must appear on >=1 slide
        ("ask exact", r"\\$11,494,200", "A1"),
        ("combined DSCR", r"1\\.49x", "A12"),
        ...
    ]
    BAN = [
        # (label, regex) — zero occurrences allowed anywhere in the deck
        ("deck-derived IRR", r"\\bIRR\\b"),          # debt register example
        ("prior-deal address", r"105 N(orth)? ?13"),  # seed remnants
        ...
    ]
    TBD_EXPECT = "slide 18 only — guarantee structure per C-X11"   # optional note

Authoring rules:
- Build MUST from the audit trail rows the blueprint actually cites (one row per
  load-bearing figure: ask, TPC, stack lines, DSCRs, NOIs, appraisal values,
  basis/LTV ratios, entitlement records, every sourced market stat).
- Build BAN from blueprint §8 verbatim (killed numerics, banned vocabulary,
  banned framings) PLUS the seed deck's deal identifiers (addresses, sponsor
  names, tenant names, market names) so no clone remnant ships.
- Sanctioned `TBD — confirm with sponsor` placeholders are NOT failures; the
  script lists every TBD location so the reviewer can confirm each is deliberate.

Exit code 0 = all checks pass; 1 = any MISSING or VIOLATION (use in the QA loop).
Keep the script output with the QA report — it is the 5b evidence.
"""
import importlib.util
import os
import re
import sys

from pptx import Presentation


def extract_slide_text(pptx_path):
    p = Presentation(pptx_path)
    slide_text = []
    for s in p.slides:
        parts = []
        for sh in s.shapes:
            if sh.has_text_frame:
                parts.append(sh.text_frame.text)
            if sh.has_table:
                for row in sh.table.rows:
                    for c in row.cells:
                        parts.append(c.text)
        slide_text.append("\n".join(parts))
    return slide_text


def main(argv=None):
    args = argv if argv is not None else sys.argv[1:]
    if len(args) != 2:
        raise SystemExit("usage: python verify-deck.py <deck.pptx> <checks.py>")
    pptx_path, checks_path = args
    spec = importlib.util.spec_from_file_location("checks", os.path.abspath(checks_path))
    checks = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(checks)
    must = getattr(checks, "MUST", [])
    ban = getattr(checks, "BAN", [])
    tbd_note = getattr(checks, "TBD_EXPECT", None)

    slide_text = extract_slide_text(pptx_path)
    fails = 0

    print("== MUST BE PRESENT (%d checks) ==" % len(must))
    for label, rx, src in must:
        hits = [i + 1 for i, t in enumerate(slide_text) if re.search(rx, t)]
        if hits:
            print("  OK  %-26s slides %s  [%s]" % (label, hits, src))
        else:
            print("  MISSING  %-22s [%s]" % (label, src))
            fails += 1

    print("== MUST BE ABSENT (%d checks) ==" % len(ban))
    for label, rx in ban:
        hits = [(i + 1, re.search(rx, t).group(0))
                for i, t in enumerate(slide_text) if re.search(rx, t)]
        if hits:
            print("  VIOLATION  %-22s %s" % (label, hits))
            fails += 1

    tbds = [(i + 1, t.count("TBD")) for i, t in enumerate(slide_text) if "TBD" in t]
    print("== TBD tokens (each must be a sanctioned, blueprint-noted placeholder) ==")
    print("  found:", tbds if tbds else "none")
    if tbd_note:
        print("  expected:", tbd_note)

    print("\nRESULT:", ("FAIL %d" % fails) if fails else "ALL CHECKS PASS")
    sys.exit(1 if fails else 0)


if __name__ == "__main__":
    main()
